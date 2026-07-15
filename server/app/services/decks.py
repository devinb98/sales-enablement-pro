"""Build a downloadable slide deck from an action plan.

Two generators, one interface:

- **Presenton** (when PRESENTON_LIVE is on and a key is set) — the hosted API
  produces a designed deck. It is credit-metered (~1 credit per slide), so it is
  off by default; development and tests never spend credits.

- **Built-in** (python-pptx) — a plain but real .pptx assembled locally. Free,
  offline, deterministic, and always available. It is the fallback when Presenton
  is disabled or fails, and it is what the tests exercise.

Either way the caller gets .pptx bytes, which are stored in Postgres and served
through an ownership-checked route. We deliberately do not hand back Presenton's
download URL: it points at a public, unauthenticated bucket, and this deck may
contain confidential deal content. What we send Presenton is the *already
synthesized plan* — summary, steps, items — not raw document passages, so less
sensitive material leaves our infrastructure than strictly necessary.
"""

import io
import logging

import requests
from flask import current_app
from pptx import Presentation
from pptx.util import Inches, Pt

log = logging.getLogger(__name__)

PRESENTON_ENDPOINT = "/api/v1/ppt/presentation/generate"
PRESENTON_TIMEOUT = 90  # synchronous generation; a few seconds per slide
PPTX_MIME = "application/vnd.openxmlformats-officedocument.presentationml.presentation"


class DeckError(RuntimeError):
    """Deck generation failed. Message is safe to surface to the user."""


def _slugify(text):
    keep = "".join(c if c.isalnum() or c in " -_" else "" for c in text)
    return "-".join(keep.split()).lower()[:60] or "deck"


def deck_filename(plan):
    company = plan.deal.company if plan.deal else "deal"
    return f"{_slugify(company)}-action-plan.pptx"


def build_content(plan):
    """The text handed to a generator. Built from the synthesized plan, not from
    the underlying documents."""
    lines = [
        f"Sales action plan for {plan.deal.company} — {plan.deal.name}.",
        "",
        "Situation:",
        plan.summary,
        "",
        "Next steps:",
    ]
    for step in plan.next_steps or []:
        lines.append(f"- {step.get('step', '')}")
    lines.append("")
    lines.append("Action items:")
    for item in plan.items:
        detail = f" — {item.detail}" if item.detail else ""
        lines.append(f"- [{item.priority}] {item.title}{detail}")
    return "\n".join(lines)


# --- Built-in generator (python-pptx) -------------------------------------


def _build_builtin_deck(plan):
    """Assemble a simple, readable .pptx locally. Returns (bytes, slide_count)."""
    prs = Presentation()
    blank = prs.slide_layouts[6]  # no placeholders; we lay out text ourselves
    slides = 0

    def add_slide(title, body_lines):
        nonlocal slides
        slide = prs.slides.add_slide(blank)
        box = slide.shapes.add_textbox(
            Inches(0.6), Inches(0.4), Inches(9), Inches(6.5)
        )
        frame = box.text_frame
        frame.word_wrap = True

        frame.text = title
        frame.paragraphs[0].font.size = Pt(30)
        frame.paragraphs[0].font.bold = True

        for line in body_lines:
            para = frame.add_paragraph()
            para.text = line
            para.font.size = Pt(16)
        slides += 1

    add_slide(
        f"{plan.deal.company} — {plan.deal.name}",
        ["Sales action plan", f"Stage: {plan.deal.stage}"],
    )
    add_slide("Where this deal stands", _wrap(plan.summary))

    steps = [s.get("step", "") for s in (plan.next_steps or [])]
    if steps:
        add_slide("Next steps", [f"• {s}" for s in steps])

    if plan.items:
        add_slide(
            "Action items",
            [f"• [{i.priority}] {i.title}" for i in plan.items],
        )

    # Sources slide — carries the deck's credibility, same as the app.
    if plan.citations:
        cites = []
        seen = set()
        for c in plan.citations:
            name = c.chunk.document.filename if c.chunk and c.chunk.document else None
            if name and name not in seen:
                seen.add(name)
                cites.append(f"• {name}")
        if cites:
            add_slide("Sources", cites)

    buffer = io.BytesIO()
    prs.save(buffer)
    return buffer.getvalue(), slides


def _wrap(text, width=90):
    """Break a paragraph into slide-friendly lines."""
    words, lines, current = text.split(), [], ""
    for word in words:
        if len(current) + len(word) + 1 > width:
            lines.append(current)
            current = word
        else:
            current = f"{current} {word}".strip()
    if current:
        lines.append(current)
    return lines


# --- Presenton generator ---------------------------------------------------


def _generate_with_presenton(plan, n_slides):
    api_key = current_app.config.get("PRESENTON_API_KEY")
    base_url = current_app.config.get("PRESENTON_BASE_URL")
    if not api_key:
        raise DeckError("Presenton API key is not configured.")

    resp = requests.post(
        f"{base_url}{PRESENTON_ENDPOINT}",
        headers={"Authorization": f"Bearer {api_key}"},
        json={
            "content": build_content(plan),
            "n_slides": n_slides,
            "language": "English",
            "template": "general",
            "export_as": "pptx",
        },
        timeout=PRESENTON_TIMEOUT,
    )
    if resp.status_code != 200:
        raise DeckError(f"Presenton returned {resp.status_code}: {resp.text[:200]}")

    download_url = resp.json().get("path")
    if not download_url:
        raise DeckError("Presenton response did not include a download path.")

    # Fetch the file server-side. The URL is public, so it must never reach the
    # browser; we store the bytes and serve them ourselves.
    file_resp = requests.get(download_url, timeout=PRESENTON_TIMEOUT)
    if file_resp.status_code != 200:
        raise DeckError("Could not download the generated deck from Presenton.")
    return file_resp.content, n_slides


# --- Public entry point ----------------------------------------------------


def generate_deck_bytes(plan, n_slides=5):
    """Return (pptx_bytes, slide_count, source).

    Uses Presenton when it is live and configured; otherwise, or if Presenton
    fails, falls back to the built-in generator so the feature always works.
    """
    if current_app.config.get("PRESENTON_LIVE") and current_app.config.get(
        "PRESENTON_API_KEY"
    ):
        try:
            data, slides = _generate_with_presenton(plan, n_slides)
            log.info("Deck for plan %d generated via Presenton", plan.id)
            return data, slides, "presenton"
        except DeckError as err:
            # A credit-metered vendor call should never be a hard failure when a
            # perfectly good local generator exists. Degrade, do not error.
            log.warning(
                "Presenton failed for plan %d, using built-in deck: %s",
                plan.id,
                err,
            )

    data, slides = _build_builtin_deck(plan)
    log.info("Deck for plan %d generated with the built-in builder", plan.id)
    return data, slides, "builtin"
