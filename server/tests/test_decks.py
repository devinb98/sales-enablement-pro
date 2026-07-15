import io

import pytest
from pptx import Presentation

from app.extensions import db as _db
from app.models import ActionItem, ActionPlan, Deal, Deck

from .conftest import login

PPTX_MAGIC = b"PK\x03\x04"  # .pptx is a zip; every zip starts with this


@pytest.fixture
def plan(user):
    deal = Deal(user_id=user.id, name="Platform Renewal", company="Acme Corp")
    _db.session.add(deal)
    _db.session.commit()
    plan = ActionPlan(
        deal_id=deal.id,
        user_id=user.id,
        summary="Acme is in discovery; a security questionnaire blocks procurement.",
        next_steps=[{"step": "Return the questionnaire by August 1.", "source_ids": [1]}],
        model_used="test",
    )
    _db.session.add(plan)
    _db.session.flush()
    _db.session.add(
        ActionItem(action_plan_id=plan.id, title="Send SOC2 report", priority="high")
    )
    _db.session.commit()
    return plan


@pytest.fixture
def rival_plan(other_user):
    deal = Deal(user_id=other_user.id, name="Rival Deal", company="Rival Inc")
    _db.session.add(deal)
    _db.session.commit()
    plan = ActionPlan(
        deal_id=deal.id, user_id=other_user.id, summary="Rival's plan.", model_used="t"
    )
    _db.session.add(plan)
    _db.session.commit()
    return plan


class TestDeckGeneration:
    def test_generates_a_real_pptx_offline(self, auth_client, plan, db):
        # PRESENTON_LIVE is off in tests, so this exercises the built-in builder —
        # no API key, no network, no credits.
        res = auth_client.post(f"/api/action-plans/{plan.id}/deck")
        assert res.status_code == 201
        body = res.get_json()
        assert body["source"] == "builtin"
        assert body["slide_count"] >= 3
        assert body["size_bytes"] > 0

        deck = db.session.get(Deck, body["id"])
        assert deck.data[:4] == PPTX_MAGIC  # it is a real Office Open XML file

    def test_the_pptx_actually_opens_and_carries_the_plan(self, auth_client, plan, db):
        deck_id = auth_client.post(f"/api/action-plans/{plan.id}/deck").get_json()["id"]
        deck = db.session.get(Deck, deck_id)

        prs = Presentation(io.BytesIO(deck.data))
        assert len(prs.slides) >= 3

        all_text = " ".join(
            shape.text
            for slide in prs.slides
            for shape in slide.shapes
            if shape.has_text_frame
        )
        # The deck reflects the plan, not a static template.
        assert "Acme Corp" in all_text
        assert "Send SOC2 report" in all_text

    def test_download_streams_the_pptx_as_an_attachment(self, auth_client, plan):
        deck_id = auth_client.post(f"/api/action-plans/{plan.id}/deck").get_json()["id"]
        res = auth_client.get(f"/api/decks/{deck_id}/download")
        assert res.status_code == 200
        assert "presentationml" in res.headers["Content-Type"]
        assert "attachment" in res.headers["Content-Disposition"]
        assert res.data[:4] == PPTX_MAGIC

    def test_download_never_exposes_a_presenton_url(self, auth_client, plan):
        """The whole reason bytes live in our database: the response is the file
        itself, never a redirect to Presenton's public bucket."""
        deck_id = auth_client.post(f"/api/action-plans/{plan.id}/deck").get_json()["id"]
        res = auth_client.get(f"/api/decks/{deck_id}/download")
        assert res.status_code == 200  # not a 302 redirect
        assert "presenton" not in res.headers.get("Location", "").lower()

    def test_list_decks_for_a_plan(self, auth_client, plan):
        auth_client.post(f"/api/action-plans/{plan.id}/deck")
        res = auth_client.get(f"/api/action-plans/{plan.id}/decks")
        assert res.status_code == 200
        assert len(res.get_json()) == 1

    def test_delete_a_deck(self, auth_client, plan, db):
        deck_id = auth_client.post(f"/api/action-plans/{plan.id}/deck").get_json()["id"]
        assert auth_client.delete(f"/api/decks/{deck_id}").status_code == 204
        assert db.session.get(Deck, deck_id) is None

    def test_deleting_a_plan_cascades_to_its_decks(self, auth_client, plan, db):
        auth_client.post(f"/api/action-plans/{plan.id}/deck")
        db.session.delete(db.session.get(ActionPlan, plan.id))
        db.session.commit()
        assert db.session.query(Deck).count() == 0


class TestDeckAuthorization:
    def test_cannot_generate_a_deck_for_another_users_plan(
        self, auth_client, rival_plan, db
    ):
        res = auth_client.post(f"/api/action-plans/{rival_plan.id}/deck")
        assert res.status_code == 404
        assert db.session.query(Deck).count() == 0

    def test_cannot_download_another_users_deck(self, client, user, rival_plan, db):
        # The rival generates a deck...
        login(client, email="rival@example.com")
        deck_id = client.post(
            f"/api/action-plans/{rival_plan.id}/deck"
        ).get_json()["id"]
        client.delete("/api/logout")

        # ...and our user tries to download it by ID.
        login(client)
        assert client.get(f"/api/decks/{deck_id}/download").status_code == 404

    def test_cannot_delete_another_users_deck(self, client, user, rival_plan, db):
        login(client, email="rival@example.com")
        deck_id = client.post(
            f"/api/action-plans/{rival_plan.id}/deck"
        ).get_json()["id"]
        client.delete("/api/logout")

        login(client)
        assert client.delete(f"/api/decks/{deck_id}").status_code == 404
        assert db.session.get(Deck, deck_id) is not None

    def test_deck_generation_requires_auth(self, client, plan):
        assert client.post(f"/api/action-plans/{plan.id}/deck").status_code == 401
